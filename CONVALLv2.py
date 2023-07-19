import collections 
import collections.abc
import pptx
import datetime
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import GetAll
presentation = Presentation('Prices.pptx')
Data = GetAll.saved_datetime
MD = Data.strftime('%d') + "." + Data.strftime('%m')
Year = Data.strftime('%Y')
print(MD)
print(Year)
for a in [0,1,2,3]:
    slide = presentation.slides[a]
    sh1 = slide.shapes[0]
    if sh1.has_text_frame:
        textpol = sh1.text_frame
    else:
        print("Форма не содержит текста")
    textpol.clear()
    p = textpol.paragraphs[0]
    run = p.add_run()
    run.text = MD
    p.alignment = PP_ALIGN.LEFT
    font = run.font
    font.name = "Montserrat SemiBolt"
    font.size = Pt(12)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)
    sh2 = slide.shapes[1]
    if sh2.has_text_frame:
        textpol1 = sh2.text_frame
    else:
        print("Форма не содержит текста")
    textpol1.clear()
    p1 = textpol1.paragraphs[0]
    run = p1.add_run()
    run.text = Year
    p1.alignment = PP_ALIGN.LEFT
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(8)
    font.color.rgb = RGBColor(255, 255, 255)
DAY0 = 0
DAY1 = 0
DAY2= 0
WEEK1 = 0
WEEK2= 0
MONTH1 = 0
MONTH2 = 0
slide = presentation.slides[0]
shape = slide.shapes[2]  
if shape.has_table:  
    table = shape.table  
else:
    print("Форма не содержит таблицы")
for i in [3, 5, 7, 9, 10, 12, 14, 15]:
    if i == 3:
        DAY0 = GetAll.PRICE1
        DAY1 = GetAll.CH_DAY1
        DAY2= GetAll.CH_DAY_PR1
        WEEK1 = GetAll.CH_W1
        WEEK2= GetAll.CH_W_PR1
        MONTH1 = GetAll.CH_M1
        MONTH2 = GetAll.CH_M_PR1
    elif i == 5:
        DAY0 = GetAll.PRICE2
        DAY1 = GetAll.CH_DAY2
        DAY2= GetAll.CH_DAY_PR2
        WEEK1 = GetAll.CH_W2
        WEEK2= GetAll.CH_W_PR2
        MONTH1 = GetAll.CH_M2
        MONTH2 = GetAll.CH_M_PR2
    elif i == 7:
        DAY0 = GetAll.PRICE3
        DAY1 = GetAll.CH_DAY3
        DAY2= GetAll.CH_DAY_PR3
        WEEK1 = GetAll.CH_W3
        WEEK2= GetAll.CH_W_PR3
        MONTH1 = GetAll.CH_M3
        MONTH2 = GetAll.CH_M_PR3
    elif i == 9:
        DAY0 = GetAll.PRICE4
        DAY1 = GetAll.CH_DAY4
        DAY2= GetAll.CH_DAY_PR4
        WEEK1 = GetAll.CH_W4
        WEEK2= GetAll.CH_W_PR4
        MONTH1 = GetAll.CH_M4
        MONTH2 = GetAll.CH_M_PR4
    elif i == 10:
        DAY0 = GetAll.PRICE5
        DAY1 = GetAll.CH_DAY5
        DAY2= GetAll.CH_DAY_PR5
        WEEK1 = GetAll.CH_W5
        WEEK2= GetAll.CH_W_PR5
        MONTH1 = GetAll.CH_M5
        MONTH2 = GetAll.CH_M_PR5
    elif i == 12:
        DAY0 = GetAll.PRICE6
        DAY1 = GetAll.CH_DAY6
        DAY2= GetAll.CH_DAY_PR6
        WEEK1 = GetAll.CH_W6
        WEEK2= GetAll.CH_W_PR6
        MONTH1 = GetAll.CH_M6
        MONTH2 = GetAll.CH_M_PR6
    elif i == 14:
        DAY0 = GetAll.PRICE7
        DAY1 = GetAll.CH_DAY7
        DAY2= GetAll.CH_DAY_PR7
        WEEK1 = GetAll.CH_W7
        WEEK2= GetAll.CH_W_PR7
        MONTH1 = GetAll.CH_M7
        MONTH2 = GetAll.CH_M_PR7
    elif i == 15:
        DAY0 = GetAll.PRICE8
        DAY1 = GetAll.CH_DAY8
        DAY2= GetAll.CH_DAY_PR8
        WEEK1 = GetAll.CH_W8
        WEEK2= GetAll.CH_W_PR8
        MONTH1 = GetAll.CH_M8
        MONTH2 = GetAll.CH_M_PR8
    else:
        print("ERROR")
    DDD0 = str (DAY0)
    DDD1 = f"{DAY1:+}"
    DDD2 = "(" + str(abs(DAY2)) + "%)"
    WWW1 = f"{WEEK1:+}"
    WWW2 = "(" + str(abs(WEEK2)) + "%)"
    MMM1 = f"{MONTH1:+}"
    MMM2 = "(" + str(abs(MONTH2)) + "%)"
    cell0 = table.cell(i, 1)
    cell0.text_frame.clear()
    cell0.text_frame.text = ""
    if DAY0 !=0:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = DDD0
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
        paragraph = cell0.text_frame.add_paragraph()
        paragraph.text = "USD/т"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(9)
        font.color.rgb = RGBColor(153, 153, 153)
    else:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
    cell1 = table.cell(i, 2)
    cell1.text_frame.clear()
    cell1.text_frame.text = ""
    if DAY1 !=0:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = DDD1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell1.text_frame.add_paragraph()
        paragraph.text = DDD2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if DAY1 == 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif DAY1 > 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell2 = table.cell(i, 3)
    cell2.text_frame.clear()
    cell2.text_frame.text = ""
    if WEEK1 != 0:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = WWW1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell2.text_frame.add_paragraph()
        paragraph.text = WWW2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if WEEK1 == 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif WEEK1 > 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell3 = table.cell(i, 4)
    cell3.text_frame.clear()
    cell3.text_frame.text = ""
    if MONTH1 != 0:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = MMM1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell3.text_frame.add_paragraph()
        paragraph.text = MMM2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if MONTH1 == 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif MONTH1 > 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 230, 231)
DAY0 = GetAll.PRICE5
DDD0 = str (DAY0)
cell4 = table.cell(10, 1)
cell4.text_frame.clear()
cell4.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell4.text_frame.add_paragraph()
    paragraph.text = "руб/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
slide = presentation.slides[1]
shape = slide.shapes[2]  
if shape.has_table:  
    table = shape.table  
else:
    print("Форма не содержит таблицы")
DAY0 = 0
DAY1 = 0
DAY2= 0
WEEK1 = 0
WEEK2= 0
MONTH1 = 0
MONTH2 = 0
for i in [3, 4, 6, 7, 8, 10, 11, 12, 14, 15]:
    if i == 3:
        DAY0 = GetAll.PRICE9
        DAY1 = GetAll.CH_DAY9
        DAY2= GetAll.CH_DAY_PR9
        WEEK1 = GetAll.CH_W9
        WEEK2= GetAll.CH_W_PR9
        MONTH1 = GetAll.CH_M9
        MONTH2 = GetAll.CH_M_PR9
    elif i == 4:
        DAY0 = GetAll.PRICE10
        DAY1 = GetAll.CH_DAY10
        DAY2= GetAll.CH_DAY_PR10
        WEEK1 = GetAll.CH_W10
        WEEK2= GetAll.CH_W_PR10
        MONTH1 = GetAll.CH_M10
        MONTH2 = GetAll.CH_M_PR10
    elif i == 6:
        DAY0 = GetAll.PRICE11
        DAY1 = GetAll.CH_DAY11
        DAY2= GetAll.CH_DAY_PR11
        WEEK1 = GetAll.CH_W11
        WEEK2= GetAll.CH_W_PR11
        MONTH1 = GetAll.CH_M11
        MONTH2 = GetAll.CH_M_PR11
    elif i == 7:
        DAY0 = GetAll.PRICE12
        DAY1 = GetAll.CH_DAY12
        DAY2= GetAll.CH_DAY_PR12
        WEEK1 = GetAll.CH_W12
        WEEK2= GetAll.CH_W_PR12
        MONTH1 = GetAll.CH_M12
        MONTH2 = GetAll.CH_M_PR12
    elif i == 8:
        DAY0 = GetAll.PRICE13
        DAY1 = GetAll.CH_DAY13
        DAY2= GetAll.CH_DAY_PR13
        WEEK1 = GetAll.CH_W13
        WEEK2= GetAll.CH_W_PR13
        MONTH1 = GetAll.CH_M13
        MONTH2 = GetAll.CH_M_PR13
    elif i == 10:
        DAY0 = GetAll.PRICE14
        DAY1 = GetAll.CH_DAY14
        DAY2= GetAll.CH_DAY_PR14
        WEEK1 = GetAll.CH_W14
        WEEK2= GetAll.CH_W_PR14
        MONTH1 = GetAll.CH_M14
        MONTH2 = GetAll.CH_M_PR14
    elif i == 11:
        DAY0 = GetAll.PRICE15
        DAY1 = GetAll.CH_DAY15
        DAY2= GetAll.CH_DAY_PR15
        WEEK1 = GetAll.CH_W15
        WEEK2= GetAll.CH_W_PR15
        MONTH1 = GetAll.CH_M15
        MONTH2 = GetAll.CH_M_PR15
    elif i == 12:
        DAY0 = GetAll.PRICE16
        DAY1 = GetAll.CH_DAY16
        DAY2= GetAll.CH_DAY_PR16
        WEEK1 = GetAll.CH_W16
        WEEK2= GetAll.CH_W_PR16
        MONTH1 = GetAll.CH_M16
        MONTH2 = GetAll.CH_M_PR16
    elif i == 14:
        DAY0 = GetAll.PRICE17
        DAY1 = GetAll.CH_DAY17
        DAY2= GetAll.CH_DAY_PR17
        WEEK1 = GetAll.CH_W17
        WEEK2= GetAll.CH_W_PR17
        MONTH1 = GetAll.CH_M17
        MONTH2 = GetAll.CH_M_PR17
    elif i == 15:
        DAY0 = GetAll.PRICE18
        DAY1 = GetAll.CH_DAY18
        DAY2= GetAll.CH_DAY_PR18
        WEEK1 = GetAll.CH_W18
        WEEK2= GetAll.CH_W_PR18
        MONTH1 = GetAll.CH_M18
        MONTH2 = GetAll.CH_M_PR18
    else:
        print("ERROR")
    DDD0 = str (DAY0)
    DDD1 = f"{DAY1:+}"
    DDD2 = "(" + str(abs(DAY2)) + "%)"
    WWW1 = f"{WEEK1:+}"
    WWW2 = "(" + str(abs(WEEK2)) + "%)"
    MMM1 = f"{MONTH1:+}"
    MMM2 = "(" + str(abs(MONTH2)) + "%)"
    cell0 = table.cell(i, 1)
    cell0.text_frame.clear()
    cell0.text_frame.text = ""
    if DAY0 !=0:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = DDD0
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
        paragraph = cell0.text_frame.add_paragraph()
        paragraph.text = "USD/т"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(9)
        font.color.rgb = RGBColor(153, 153, 153)
    else:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
    cell1 = table.cell(i, 2)
    cell1.text_frame.clear()
    cell1.text_frame.text = ""
    if DAY1 !=0:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = DDD1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell1.text_frame.add_paragraph()
        paragraph.text = DDD2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if DAY1 == 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif DAY1 > 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell2 = table.cell(i, 3)
    cell2.text_frame.clear()
    cell2.text_frame.text = ""
    if WEEK1 != 0:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = WWW1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell2.text_frame.add_paragraph()
        paragraph.text = WWW2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if WEEK1 == 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif WEEK1 > 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell3 = table.cell(i, 4)
    cell3.text_frame.clear()
    cell3.text_frame.text = ""
    if MONTH1 != 0:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = MMM1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell3.text_frame.add_paragraph()
        paragraph.text = MMM2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if MONTH1 == 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif MONTH1 > 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 230, 231)
DAY0 = GetAll.PRICE13
DDD0 = str (DAY0)
cell4 = table.cell(8, 1)
cell4.text_frame.clear()
cell4.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell4.text_frame.add_paragraph()
    paragraph.text = "руб/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
DAY0 = GetAll.PRICE16
DDD0 = str (DAY0)
cell5 = table.cell(12, 1)
cell5.text_frame.clear()
cell5.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell5.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell5.text_frame.add_paragraph()
    paragraph.text = "руб/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell5.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
DAY0 = GetAll.PRICE18
DDD0 = str (DAY0)
cell6 = table.cell(15, 1)
cell6.text_frame.clear()
cell6.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell6.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell6.text_frame.add_paragraph()
    paragraph.text = "руб/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell6.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
slide = presentation.slides[2]
shape = slide.shapes[2]  
if shape.has_table:  
    table = shape.table  
else:
    print("Форма не содержит таблицы")
DAY0 = 0
DAY1 = 0
DAY2= 0
WEEK1 = 0
WEEK2= 0
MONTH1 = 0
MONTH2 = 0
for i in [3, 4, 5, 6, 8, 9, 10, 11, 12, 14]:
    if i == 3:
        DAY0 = GetAll.PRICE19
        DAY1 = GetAll.CH_DAY19
        DAY2= GetAll.CH_DAY_PR19
        WEEK1 = GetAll.CH_W19
        WEEK2= GetAll.CH_W_PR19
        MONTH1 = GetAll.CH_M19
        MONTH2 = GetAll.CH_M_PR19
    elif i == 4:
        DAY0 = GetAll.PRICE20
        DAY1 = GetAll.CH_DAY20
        DAY2= GetAll.CH_DAY_PR20
        WEEK1 = GetAll.CH_W20
        WEEK2= GetAll.CH_W_PR20
        MONTH1 = GetAll.CH_M20
        MONTH2 = GetAll.CH_M_PR20
    elif i == 5:
        DAY0 = GetAll.PRICE21
        DAY1 = GetAll.CH_DAY21
        DAY2= GetAll.CH_DAY_PR21
        WEEK1 = GetAll.CH_W21
        WEEK2= GetAll.CH_W_PR21
        MONTH1 = GetAll.CH_M21
        MONTH2 = GetAll.CH_M_PR21
    elif i == 6:
        DAY0 = GetAll.PRICE22
        DAY1 = GetAll.CH_DAY22
        DAY2= GetAll.CH_DAY_PR22
        WEEK1 = GetAll.CH_W22
        WEEK2= GetAll.CH_W_PR22
        MONTH1 = GetAll.CH_M22
        MONTH2 = GetAll.CH_M_PR22
    elif i == 8:
        DAY0 = GetAll.PRICE23
        DAY1 = GetAll.CH_DAY23
        DAY2= GetAll.CH_DAY_PR23
        WEEK1 = GetAll.CH_W23
        WEEK2= GetAll.CH_W_PR23
        MONTH1 = GetAll.CH_M23
        MONTH2 = GetAll.CH_M_PR23
    elif i == 9:
        DAY0 = GetAll.PRICE24
        DAY1 = GetAll.CH_DAY24
        DAY2= GetAll.CH_DAY_PR24
        WEEK1 = GetAll.CH_W24
        WEEK2= GetAll.CH_W_PR24
        MONTH1 = GetAll.CH_M24
        MONTH2 = GetAll.CH_M_PR24
    elif i == 10:
        DAY0 = GetAll.PRICE25
        DAY1 = GetAll.CH_DAY25
        DAY2= GetAll.CH_DAY_PR25
        WEEK1 = GetAll.CH_W25
        WEEK2= GetAll.CH_W_PR25
        MONTH1 = GetAll.CH_M25
        MONTH2 = GetAll.CH_M_PR25
    elif i == 11:
        DAY0 = GetAll.PRICE26
        DAY1 = GetAll.CH_DAY26
        DAY2= GetAll.CH_DAY_PR26
        WEEK1 = GetAll.CH_W26
        WEEK2= GetAll.CH_W_PR26
        MONTH1 = GetAll.CH_M26
        MONTH2 = GetAll.CH_M_PR26
    elif i == 12:
        DAY0 = GetAll.PRICE27
        DAY1 = GetAll.CH_DAY27
        DAY2= GetAll.CH_DAY_PR27
        WEEK1 = GetAll.CH_W27
        WEEK2= GetAll.CH_W_PR27
        MONTH1 = GetAll.CH_M27
        MONTH2 = GetAll.CH_M_PR27
    elif i == 14:
        DAY0 = GetAll.PRICE28
        DAY1 = GetAll.CH_DAY28
        DAY2= GetAll.CH_DAY_PR28
        WEEK1 = GetAll.CH_W28
        WEEK2= GetAll.CH_W_PR28
        MONTH1 = GetAll.CH_M28
        MONTH2 = GetAll.CH_M_PR28
    else:
        print("ERROR")
    DDD0 = str (DAY0)
    DDD1 = f"{DAY1:+}"
    DDD2 = "(" + str(abs(DAY2)) + "%)"
    WWW1 = f"{WEEK1:+}"
    WWW2 = "(" + str(abs(WEEK2)) + "%)"
    MMM1 = f"{MONTH1:+}"
    MMM2 = "(" + str(abs(MONTH2)) + "%)"
    cell0 = table.cell(i, 1)
    cell0.text_frame.clear()
    cell0.text_frame.text = ""
    if DAY0 !=0:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = DDD0
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
        paragraph = cell0.text_frame.add_paragraph()
        paragraph.text = "USD/т"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(9)
        font.color.rgb = RGBColor(153, 153, 153)
    else:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
    cell1 = table.cell(i, 2)
    cell1.text_frame.clear()
    cell1.text_frame.text = ""
    if DAY1 !=0:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = DDD1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell1.text_frame.add_paragraph()
        paragraph.text = DDD2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if DAY1 == 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif DAY1 > 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell2 = table.cell(i, 3)
    cell2.text_frame.clear()
    cell2.text_frame.text = ""
    if WEEK1 != 0:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = WWW1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell2.text_frame.add_paragraph()
        paragraph.text = WWW2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if WEEK1 == 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif WEEK1 > 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell3 = table.cell(i, 4)
    cell3.text_frame.clear()
    cell3.text_frame.text = ""
    if MONTH1 != 0:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = MMM1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell3.text_frame.add_paragraph()
        paragraph.text = MMM2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if MONTH1 == 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif MONTH1 > 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 230, 231)
DAY0 = GetAll.PRICE19
DDD0 = str (DAY0)
cell4 = table.cell(3, 1)
cell4.text_frame.clear()
cell4.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell4.text_frame.add_paragraph()
    paragraph.text = "CNY/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
DAY0 = GetAll.PRICE23
DDD0 = str (DAY0)
cell5 = table.cell(8, 1)
cell5.text_frame.clear()
cell5.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell5.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell5.text_frame.add_paragraph()
    paragraph.text = "CNY/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell5.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
DAY0 = GetAll.PRICE28
DDD0 = str (DAY0)
cell6 = table.cell(14, 1)
cell6.text_frame.clear()
cell6.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell6.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell6.text_frame.add_paragraph()
    paragraph.text = "USD/1% Mn в смт"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell6.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
slide = presentation.slides[3]
shape = slide.shapes[2]  
if shape.has_table:  
    table = shape.table  
else:
    print("Форма не содержит таблицы")
DAY0 = 0
DAY1 = 0
DAY2= 0
WEEK1 = 0
WEEK2= 0
MONTH1 = 0
MONTH2 = 0
for i in [3, 4, 5, 6, 7, 9, 10, 11, 13]:
    if i == 3:
        DAY0 = GetAll.PRICE29
        DAY1 = GetAll.CH_DAY29
        DAY2= GetAll.CH_DAY_PR29
        WEEK1 = GetAll.CH_W29
        WEEK2= GetAll.CH_W_PR29
        MONTH1 = GetAll.CH_M29
        MONTH2 = GetAll.CH_M_PR29
    elif i == 4:
        DAY0 = GetAll.PRICE30
        DAY1 = GetAll.CH_DAY30
        DAY2= GetAll.CH_DAY_PR30
        WEEK1 = GetAll.CH_W30
        WEEK2= GetAll.CH_W_PR30
        MONTH1 = GetAll.CH_M30
        MONTH2 = GetAll.CH_M_PR30
    elif i == 5:
        DAY0 = GetAll.PRICE31
        DAY1 = GetAll.CH_DAY31
        DAY2= GetAll.CH_DAY_PR31
        WEEK1 = GetAll.CH_W31
        WEEK2= GetAll.CH_W_PR31
        MONTH1 = GetAll.CH_M31
        MONTH2 = GetAll.CH_M_PR31
    elif i == 6:
        DAY0 = GetAll.PRICE32
        DAY1 = GetAll.CH_DAY32
        DAY2= GetAll.CH_DAY_PR32
        WEEK1 = GetAll.CH_W32
        WEEK2= GetAll.CH_W_PR32
        MONTH1 = GetAll.CH_M32
        MONTH2 = GetAll.CH_M_PR32
    elif i == 7:
        DAY0 = GetAll.PRICE33
        DAY1 = GetAll.CH_DAY33
        DAY2= GetAll.CH_DAY_PR33
        WEEK1 = GetAll.CH_W33
        WEEK2= GetAll.CH_W_PR33
        MONTH1 = GetAll.CH_M33
        MONTH2 = GetAll.CH_M_PR33
    elif i == 9:
        DAY0 = GetAll.PRICE34
        DAY1 = GetAll.CH_DAY34
        DAY2= GetAll.CH_DAY_PR34
        WEEK1 = GetAll.CH_W34
        WEEK2= GetAll.CH_W_PR34
        MONTH1 = GetAll.CH_M34
        MONTH2 = GetAll.CH_M_PR34
    elif i == 10:
        DAY0 = GetAll.PRICE35
        DAY1 = GetAll.CH_DAY35
        DAY2= GetAll.CH_DAY_PR35
        WEEK1 = GetAll.CH_W35
        WEEK2= GetAll.CH_W_PR35
        MONTH1 = GetAll.CH_M35
        MONTH2 = GetAll.CH_M_PR35
    elif i == 11:
        DAY0 = GetAll.PRICE36
        DAY1 = GetAll.CH_DAY36
        DAY2= GetAll.CH_DAY_PR36
        WEEK1 = GetAll.CH_W36
        WEEK2= GetAll.CH_W_PR36
        MONTH1 = GetAll.CH_M36
        MONTH2 = GetAll.CH_M_PR36
    elif i == 13:
        DAY0 = GetAll.PRICE37
        DAY1 = GetAll.CH_DAY37
        DAY2= GetAll.CH_DAY_PR37
        WEEK1 = GetAll.CH_W37
        WEEK2= GetAll.CH_W_PR37
        MONTH1 = GetAll.CH_M37
        MONTH2 = GetAll.CH_M_PR37
    else:
        print("ERROR")
    DDD0 = str (DAY0)
    DDD1 = f"{DAY1:+}"
    DDD2 = "(" + str(abs(DAY2)) + "%)"
    WWW1 = f"{WEEK1:+}"
    WWW2 = "(" + str(abs(WEEK2)) + "%)"
    MMM1 = f"{MONTH1:+}"
    MMM2 = "(" + str(abs(MONTH2)) + "%)"
    cell0 = table.cell(i, 1)
    cell0.text_frame.clear()
    cell0.text_frame.text = ""
    if DAY0 !=0:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = DDD0
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
        paragraph = cell0.text_frame.add_paragraph()
        paragraph.text = "USDc/фунт Cr"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(9)
        font.color.rgb = RGBColor(153, 153, 153)
    else:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
    cell1 = table.cell(i, 2)
    cell1.text_frame.clear()
    cell1.text_frame.text = ""
    if DAY1 !=0:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = DDD1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell1.text_frame.add_paragraph()
        paragraph.text = DDD2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if DAY1 == 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif DAY1 > 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell2 = table.cell(i, 3)
    cell2.text_frame.clear()
    cell2.text_frame.text = ""
    if WEEK1 != 0:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = WWW1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell2.text_frame.add_paragraph()
        paragraph.text = WWW2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if WEEK1 == 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif WEEK1 > 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell3 = table.cell(i, 4)
    cell3.text_frame.clear()
    cell3.text_frame.text = ""
    if MONTH1 != 0:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = MMM1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell3.text_frame.add_paragraph()
        paragraph.text = MMM2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if MONTH1 == 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif MONTH1 > 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 230, 231)
DAY0 = GetAll.PRICE29
DDD0 = str (DAY0)
cell4 = table.cell(3, 1)
cell4.text_frame.clear()
cell4.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell4.text_frame.add_paragraph()
    paragraph.text = "CNY/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
DAY0 = GetAll.PRICE34
DDD0 = str (DAY0)
cell5 = table.cell(9, 1)
cell5.text_frame.clear()
cell5.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell5.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell5.text_frame.add_paragraph()
    paragraph.text = "CNY/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell5.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
DAY0 = GetAll.PRICE37
DDD0 = str (DAY0)
cell6 = table.cell(13, 1)
cell6.text_frame.clear()
cell6.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell6.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell6.text_frame.add_paragraph()
    paragraph.text = "USD/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell6.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
presentation.save('updated_example.pptx')
print("ЗАПИСЬ УСПЕШНА")
import os
import win32com.client as win32
import comtypes
from pdf2image import convert_from_path
comtypes.CoInitialize()
input_path = r'C:\Users\Марк\Desktop\tgbot\updated_example.pptx'
output_path = r'C:\Users\Марк\Desktop\tgbot\super.pdf'
powerpoint = win32.Dispatch('Powerpoint.Application')
presentation = powerpoint.Presentations.Open(input_path)
presentation.SaveAs(output_path , 32)
presentation.Close()
powerpoint.Quit()
pdf_path = 'super.pdf'
images = convert_from_path(pdf_path, first_page=1, last_page=4)
for i, image in enumerate(images):
    image.save(f'page{i+1}.png', 'PNG')