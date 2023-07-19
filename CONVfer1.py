import collections 
import collections.abc
import pptx
import datetime
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import Ferro1Excel
presentation = Presentation('Prices.pptx')
slide = presentation.slides[2]
Data = Ferro1Excel.DATA
MD = Data.strftime('%d') + "." + Data.strftime('%m')
Year = Data.strftime('%Y')
print(MD)
print(Year)
sh1 = slide.shapes[0]
if sh1.has_text_frame:
    textpol = sh1.text_frame
else:
    print("Форма не содержит текста")
textpol.clear()
p = textpol.paragraphs[0]
run = p.add_run()
run.text = MD
p.alignment = PP_ALIGN.LEFT
font = run.font
font.name = "Montserrat SemiBolt"
font.size = Pt(12)
font.bold = True
font.color.rgb = RGBColor(255, 255, 255)
sh2 = slide.shapes[1]
if sh2.has_text_frame:
    textpol1 = sh2.text_frame
else:
    print("Форма не содержит текста")
textpol1.clear()
p1 = textpol1.paragraphs[0]
run = p1.add_run()
run.text = Year
p1.alignment = PP_ALIGN.LEFT
font = run.font
font.name = "Montserrat"
font.size = Pt(8)
font.color.rgb = RGBColor(255, 255, 255)
shape = slide.shapes[2]  
if shape.has_table:  
    table = shape.table  
else:
    print("Форма не содержит таблицы")
DAY0 = 0
DAY1 = 0
DAY2= 0
WEEK1 = 0
WEEK2= 0
MONTH1 = 0
MONTH2 = 0
for i in [3, 4, 5, 6, 8, 9, 10, 11, 12, 14]:
    if i == 3:
        DAY0 = Ferro1Excel.PRICE1
        DAY1 = Ferro1Excel.CH_DAY1
        DAY2= Ferro1Excel.CH_DAY_PR1
        WEEK1 = Ferro1Excel.CH_W1
        WEEK2= Ferro1Excel.CH_W_PR1
        MONTH1 = Ferro1Excel.CH_M1
        MONTH2 = Ferro1Excel.CH_M_PR1
    elif i == 4:
        DAY0 = Ferro1Excel.PRICE2
        DAY1 = Ferro1Excel.CH_DAY2
        DAY2= Ferro1Excel.CH_DAY_PR2
        WEEK1 = Ferro1Excel.CH_W2
        WEEK2= Ferro1Excel.CH_W_PR2
        MONTH1 = Ferro1Excel.CH_M2
        MONTH2 = Ferro1Excel.CH_M_PR2
    elif i == 5:
        DAY0 = Ferro1Excel.PRICE3
        DAY1 = Ferro1Excel.CH_DAY3
        DAY2= Ferro1Excel.CH_DAY_PR3
        WEEK1 = Ferro1Excel.CH_W3
        WEEK2= Ferro1Excel.CH_W_PR3
        MONTH1 = Ferro1Excel.CH_M3
        MONTH2 = Ferro1Excel.CH_M_PR3
    elif i == 6:
        DAY0 = Ferro1Excel.PRICE4
        DAY1 = Ferro1Excel.CH_DAY4
        DAY2= Ferro1Excel.CH_DAY_PR4
        WEEK1 = Ferro1Excel.CH_W4
        WEEK2= Ferro1Excel.CH_W_PR4
        MONTH1 = Ferro1Excel.CH_M4
        MONTH2 = Ferro1Excel.CH_M_PR4
    elif i == 8:
        DAY0 = Ferro1Excel.PRICE5
        DAY1 = Ferro1Excel.CH_DAY5
        DAY2= Ferro1Excel.CH_DAY_PR5
        WEEK1 = Ferro1Excel.CH_W5
        WEEK2= Ferro1Excel.CH_W_PR5
        MONTH1 = Ferro1Excel.CH_M5
        MONTH2 = Ferro1Excel.CH_M_PR5
    elif i == 9:
        DAY0 = Ferro1Excel.PRICE6
        DAY1 = Ferro1Excel.CH_DAY6
        DAY2= Ferro1Excel.CH_DAY_PR6
        WEEK1 = Ferro1Excel.CH_W6
        WEEK2= Ferro1Excel.CH_W_PR6
        MONTH1 = Ferro1Excel.CH_M6
        MONTH2 = Ferro1Excel.CH_M_PR6
    elif i == 10:
        DAY0 = Ferro1Excel.PRICE7
        DAY1 = Ferro1Excel.CH_DAY7
        DAY2= Ferro1Excel.CH_DAY_PR7
        WEEK1 = Ferro1Excel.CH_W7
        WEEK2= Ferro1Excel.CH_W_PR7
        MONTH1 = Ferro1Excel.CH_M7
        MONTH2 = Ferro1Excel.CH_M_PR7
    elif i == 11:
        DAY0 = Ferro1Excel.PRICE8
        DAY1 = Ferro1Excel.CH_DAY8
        DAY2= Ferro1Excel.CH_DAY_PR8
        WEEK1 = Ferro1Excel.CH_W8
        WEEK2= Ferro1Excel.CH_W_PR8
        MONTH1 = Ferro1Excel.CH_M8
        MONTH2 = Ferro1Excel.CH_M_PR8
    elif i == 12:
        DAY0 = Ferro1Excel.PRICE9
        DAY1 = Ferro1Excel.CH_DAY9
        DAY2= Ferro1Excel.CH_DAY_PR9
        WEEK1 = Ferro1Excel.CH_W9
        WEEK2= Ferro1Excel.CH_W_PR9
        MONTH1 = Ferro1Excel.CH_M9
        MONTH2 = Ferro1Excel.CH_M_PR9
    elif i == 14:
        DAY0 = Ferro1Excel.PRICE10
        DAY1 = Ferro1Excel.CH_DAY10
        DAY2= Ferro1Excel.CH_DAY_PR10
        WEEK1 = Ferro1Excel.CH_W10
        WEEK2= Ferro1Excel.CH_W_PR10
        MONTH1 = Ferro1Excel.CH_M10
        MONTH2 = Ferro1Excel.CH_M_PR10
    else:
        print("ERROR")
    DDD0 = str (DAY0)
    DDD1 = f"{DAY1:+}"
    DDD2 = "(" + str(abs(DAY2)) + "%)"
    WWW1 = f"{WEEK1:+}"
    WWW2 = "(" + str(abs(WEEK2)) + "%)"
    MMM1 = f"{MONTH1:+}"
    MMM2 = "(" + str(abs(MONTH2)) + "%)"
    cell0 = table.cell(i, 1)
    cell0.text_frame.clear()
    cell0.text_frame.text = ""
    if DAY0 !=0:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = DDD0
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
        paragraph = cell0.text_frame.add_paragraph()
        paragraph.text = "USD/т"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(9)
        font.color.rgb = RGBColor(153, 153, 153)
    else:
        paragraph = cell0.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat"
        font.size = Pt(11)
        font.bold = True
        font.color.rgb = RGBColor(77, 77, 77)
    cell1 = table.cell(i, 2)
    cell1.text_frame.clear()
    cell1.text_frame.text = ""
    if DAY1 !=0:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = DDD1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell1.text_frame.add_paragraph()
        paragraph.text = DDD2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell1.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if DAY1 == 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif DAY1 > 0:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell2 = table.cell(i, 3)
    cell2.text_frame.clear()
    cell2.text_frame.text = ""
    if WEEK1 != 0:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = WWW1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell2.text_frame.add_paragraph()
        paragraph.text = WWW2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell2.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if WEEK1 == 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif WEEK1 > 0:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(255, 230, 231)
    cell3 = table.cell(i, 4)
    cell3.text_frame.clear()
    cell3.text_frame.text = ""
    if MONTH1 != 0:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = MMM1
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
        paragraph = cell3.text_frame.add_paragraph()
        paragraph.text = MMM2
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat Light"
        font.size = Pt(9)
    else:
        paragraph = cell3.text_frame.paragraphs[0]
        paragraph.text = "-"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = paragraph.runs[0]
        font = run.font
        font.name = "Montserrat SemiBold"
        font.size = Pt(11)
        font.bold = True
    if MONTH1 == 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 255, 255)
    elif MONTH1 > 0:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(231, 247, 240)
    else:
        cell3.fill.solid()
        cell3.fill.fore_color.rgb = RGBColor(255, 230, 231)
DAY0 = Ferro1Excel.PRICE1
DDD0 = str (DAY0)
cell4 = table.cell(3, 1)
cell4.text_frame.clear()
cell4.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell4.text_frame.add_paragraph()
    paragraph.text = "CNY/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell4.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
DAY0 = Ferro1Excel.PRICE5
DDD0 = str (DAY0)
cell5 = table.cell(8, 1)
cell5.text_frame.clear()
cell5.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell5.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell5.text_frame.add_paragraph()
    paragraph.text = "CNY/т"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell5.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
DAY0 = Ferro1Excel.PRICE10
DDD0 = str (DAY0)
cell6 = table.cell(14, 1)
cell6.text_frame.clear()
cell6.text_frame.text = ""
if DAY0 !=0:
    paragraph = cell6.text_frame.paragraphs[0]
    paragraph.text = DDD0
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
    paragraph = cell6.text_frame.add_paragraph()
    paragraph.text = "USD/1% Mn в смт"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(9)
    font.color.rgb = RGBColor(153, 153, 153)
else:
    paragraph = cell6.text_frame.paragraphs[0]
    paragraph.text = "-"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = paragraph.runs[0]
    font = run.font
    font.name = "Montserrat"
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = RGBColor(77, 77, 77)
presentation.save('updated_example.pptx')
print("ЗАПИСЬ УСПЕШНА")
import os
import win32com.client as win32
import comtypes
from pdf2image import convert_from_path
comtypes.CoInitialize()
input_path = r'C:\Users\Марк\Desktop\tgbot\updated_example.pptx'
output_path = r'C:\Users\Марк\Desktop\tgbot\super.pdf'
powerpoint = win32.Dispatch('Powerpoint.Application')
presentation = powerpoint.Presentations.Open(input_path)
presentation.SaveAs(output_path , 32)
presentation.Close()
powerpoint.Quit()
pdf_path = 'super.pdf'
images = convert_from_path(pdf_path, first_page=3, last_page=3)
images[0].save('page3.png', 'PNG')